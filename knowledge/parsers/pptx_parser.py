
"""
PowerPoint 文档解析器
"""
from pathlib import Path
from typing import Optional
from .base import BaseParser

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class PptxParser(BaseParser):
    """PowerPoint (.pptx) 文档解析器"""
    
    def supported_extensions(self) -> list[str]:
        return [".pptx"]
    
    def to_markdown(self, file_path: str, output_path: Optional[str] = None) -> str:
        if not HAS_PPTX:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")
        
        prs = Presentation(file_path)
        markdown_parts = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts = [f"---\n\n## 幻灯片 {slide_num}"]
            
            # 提取标题和内容
            title_text = ""
            content_texts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # 判断是否是标题
                    if shape.is_placeholder and shape.placeholder_format.type == 1:
                        title_text = shape.text.strip()
                    else:
                        content_texts.append(shape.text.strip())
            
            if title_text:
                slide_parts.append(f"\n# {title_text}")
            
            if content_texts:
                slide_parts.append("\n" + "\n\n".join(content_texts))
            
            # 提取备注
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                slide_parts.append(f"\n> **备注**: {notes_text}")
            
            if len(slide_parts) > 1:
                markdown_parts.append("\n".join(slide_parts))
        
        markdown_text = "\n\n".join(markdown_parts)
        
        if output_path:
            self._save_markdown(markdown_text, output_path)
        
        return markdown_text

